import os
import fitz  # PyMuPDF
import yaml
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO


def load_yaml(path):
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def merge_pdfs(pdf_paths, output_path):
    merged = fitz.open()

    for path in pdf_paths:
        if not isinstance(path, str) or not path.strip():
            continue

        if not os.path.isfile(path):
            raise FileNotFoundError(f"Missing file: {path}")

        doc = fitz.open(path)
        merged.insert_pdf(doc)
        doc.close()

    merged.save(output_path)
    merged.close()


def add_title_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        prs.slide_width - Inches(2),
        Inches(2)
    )

    tf = textbox.text_frame
    tf.text = title_text

    return slide


def add_pdf_pages(prs, pdf_path, zoom=3):
    doc = fitz.open(pdf_path)

    for page in doc:
        pix = page.get_pixmap(
            matrix=fitz.Matrix(zoom, zoom),
            colorspace=fitz.csRGB,
            alpha=False
        )

        img_bytes = pix.tobytes("jpg", jpg_quality=80)
        image_stream = BytesIO(img_bytes)

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        slide.shapes.add_picture(
            image_stream,
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()


def pdfs_to_pptx_with_titles(pdf_paths, output_path, zoom=3):
    prs = Presentation()

    for path in pdf_paths:
        if not isinstance(path, str) or not path.strip():
            continue

        if not os.path.isfile(path):
            raise FileNotFoundError(f"Missing file: {path}")

        filename = os.path.basename(path)

        # Title slide
        add_title_slide(prs, filename)

        # Content slides
        add_pdf_pages(prs, path, zoom=zoom)

    prs.save(output_path)


def main():
    config = load_yaml("order.yaml")

    base_out = "output"
    pdf_out = os.path.join(base_out, "merged_pdfs")
    pptx_out = os.path.join(base_out, "output_pptx_import")

    os.makedirs(pdf_out, exist_ok=True)
    os.makedirs(pptx_out, exist_ok=True)

    for key, item in config["pdfs"].items():
        name = item["name"]
        pdf_paths = item["pdfs_merge"]

        # default: process = True
        should_process = item.get("process", True)

        if not should_process:
            print(f"\nSkipping: {name}")
            continue

        print(f"\nProcessing: {name}")

        # Merge PDFs
        merged_pdf_path = os.path.join(pdf_out, f"{name}.pdf")
        merge_pdfs(pdf_paths, merged_pdf_path)
        print(f"Saved PDF: {merged_pdf_path}")

        # Create PPTX with title slides
        pptx_path = os.path.join(pptx_out, f"{name}.pptx")
        print("Creating PPTX with title slides...")

        pdfs_to_pptx_with_titles(pdf_paths, pptx_path, zoom=3)

        print(f"Saved PPTX: {pptx_path}")


if __name__ == "__main__":
    main()